
from pptx import Presentation
import copy

def merge_pptx(
        input_files: list,
        res_file: str,
        tmpl_pptx: str = None,
        tmpl_index: int = 0) -> any:
    
    assert len(input_files) > 1
    merged_pptx = None
    tmpl_slide_layout = None
    if tmpl_pptx:
        # if has tmpl_pptx
        merged_pptx = Presentation(tmpl_pptx)
        tmpl_slide_layout = merged_pptx.slide_layouts[tmpl_index]
    for file in input_files:
        if merged_pptx is None:
            merged_pptx = Presentation(file)
            tmpl_slide_layout = merged_pptx.slide_layouts[len(merged_pptx.slides)]
            continue
        present = Presentation(file)
        for slide in present.slides:
            tmp_slide = merged_pptx.slides.add_slide(tmpl_slide_layout)
                # tmp_slide = merged_pptx.slides.add_slide(present.slide_layouts[i])
            # remove default elements
            for shp in tmp_slide.shapes:
                tmp_slide.shapes.element.remove(shp.element)
            # copy elements
            for shp in slide.shapes:
                el = shp.element
                new_el = copy.deepcopy(el)
                tmp_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
    merged_pptx.save(res_file)
    # return merged_pptx

help = """-d | dir path
-i | input filenames
-o | output filename
-t | templete file | not yet
-ti | templete index | not yet
"""

if __name__ == "__main__":
    # terminal command line
    import sys, os
    argv = sys.argv
    input_files = []
    output_file = None
    if '-d' in argv:
        value_index = argv.index('-d') + 1
        if len(argv) > value_index:
            input_files = [os.path.join(argv[value_index], it) for it in os.listdir(argv[value_index]) if it.lower().endswith('.pptx')]
    elif '-i' in argv:
        value_index = argv.index('-i') + 1
        if len(argv) > value_index:
            files = argv[value_index:]
            end_index = files.index('-')
            files = files[:end_index]
            input_files = [it for it in files if it.lower().endswith('pptx')]
    if '-o' in argv:
        value_index = argv.index('-o') + 1
        output_file = argv[value_index] if len(argv) > value_index else None
    if input_files and output_file:
        # 유효
        merge_pptx(input_files, output_file)
    else:
        print(help)